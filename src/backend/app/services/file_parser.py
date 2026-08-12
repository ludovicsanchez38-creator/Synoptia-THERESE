"""
THÉRÈSE v2 - File Parser Service

Extracts text content from various file types for indexing.
"""

import logging
import mimetypes
from pathlib import Path
from typing import Generator

logger = logging.getLogger(__name__)

MAX_FILE_SIZE_BYTES = 50 * 1024 * 1024  # 50 Mo
MAX_PDF_PAGES = 100
# Inventaire des capacités du 13/08/2026 : le PDF était plafonné à 100 pages et
# le CSV à 500 lignes, chacun avec une mention de troncature. Le tableur, lui,
# parcourait TOUTES les feuilles et TOUTES les lignes sans borne — un classeur
# de plusieurs dizaines de milliers de lignes produisait un texte énorme, lent
# à découper et coûteux à envoyer, sans que rien ne le signale.
MAX_XLSX_LIGNES = 2000


# Supported extensions
#
# Inventaire des capacités du 13/08/2026 : six extensions étaient acceptées à
# l'indexation sans figurer ici. Le contrôle d'entrée les laissait passer, le
# parseur ne savait pas quoi en faire, et l'indexation « réussissait » sur un
# contenu vide sans que rien ne le signale. Ce sont pourtant des fichiers texte,
# que lire ne demande aucune bibliothèque — seulement de les déclarer.
TEXT_EXTENSIONS = {
    ".txt", ".md", ".markdown", ".rst", ".log",
    ".cfg", ".conf", ".ini", ".org", ".tex",
}
CODE_EXTENSIONS = {
    ".scala",
    ".py",
    ".js",
    ".ts",
    ".tsx",
    ".jsx",
    ".html",
    ".css",
    ".scss",
    ".json",
    ".yaml",
    ".yml",
    ".toml",
    ".xml",
    ".sql",
    ".sh",
    ".bash",
    ".zsh",
    ".fish",
    ".r",
    ".rs",
    ".go",
    ".java",
    ".c",
    ".cpp",
    ".h",
    ".hpp",
    ".swift",
    ".kt",
    ".rb",
    ".php",
    ".lua",
    ".vim",
    ".el",
}
PRESENTATION_EXTENSIONS = {".pptx"}
CSV_EXTENSIONS = {".csv", ".tsv"}


def extract_text(file_path: Path) -> str | None:
    """
    Extract text content from a file.

    Args:
        file_path: Path to the file

    Returns:
        Extracted text or None if unsupported/error
    """
    if not file_path.exists():
        logger.warning(f"File not found: {file_path}")
        return None

    # Vérification de la taille maximale (SEC-016, PERF-009)
    file_size = file_path.stat().st_size
    if file_size > MAX_FILE_SIZE_BYTES:
        raise ValueError("Le fichier dépasse la limite de 50 Mo")

    ext = file_path.suffix.lower()

    try:
        # Text and code files
        if ext in TEXT_EXTENSIONS or ext in CODE_EXTENSIONS:
            return _extract_plain_text(file_path)

        # CSV files
        if ext in CSV_EXTENSIONS:
            return _extract_csv(file_path)

        # PDF files
        if ext == ".pdf":
            return _extract_pdf(file_path)

        # Word documents
        if ext in {".docx", ".doc"}:
            return _extract_docx(file_path)

        # Excel spreadsheets
        if ext == ".xlsx":
            return _extract_xlsx(file_path)

        # Présentations
        if ext in PRESENTATION_EXTENSIONS:
            return _extract_pptx(file_path)

        # Unsupported
        logger.warning(f"Unsupported file type: {ext}")
        return None

    except Exception as e:
        logger.error(f"Error extracting text from {file_path}: {e}")
        return None


def _extract_plain_text(file_path: Path) -> str:
    """Extract text from plain text files."""
    # Try common encodings
    encodings = ["utf-8", "latin-1", "cp1252"]

    for encoding in encodings:
        try:
            return file_path.read_text(encoding=encoding)
        except UnicodeDecodeError:
            continue

    raise ValueError(f"Could not decode file with any of {encodings}")


def _extract_csv(file_path: Path) -> str:
    """Extract text from CSV files."""
    import csv
    from io import StringIO

    text = _extract_plain_text(file_path)
    reader = csv.reader(StringIO(text))

    # Convert to readable format
    lines = []
    for i, row in enumerate(reader):
        if i == 0:
            # Header
            lines.append(" | ".join(row))
            lines.append("-" * 40)
        else:
            lines.append(" | ".join(row))

        # Limit rows for very large CSVs
        if i > 500:
            lines.append(f"... ({i}+ rows total)")
            break

    return "\n".join(lines)


def _extract_pdf(file_path: Path) -> str:
    """Extract text from PDF files using pypdf."""
    try:
        from pypdf import PdfReader
    except ImportError:
        logger.warning("pypdf not installed, cannot extract PDF")
        return "[PDF extraction unavailable - install pypdf]"

    reader = PdfReader(file_path)
    total_pages = len(reader.pages)
    pages_to_process = min(total_pages, MAX_PDF_PAGES)
    text_parts = []

    for i in range(pages_to_process):
        try:
            text = reader.pages[i].extract_text()
            if text:
                text_parts.append(f"--- Page {i + 1} ---\n{text}")
        except Exception as e:
            logger.warning(f"Error extracting page {i + 1}: {e}")

    result = "\n\n".join(text_parts)

    if total_pages > MAX_PDF_PAGES:
        logger.info(f"PDF tronqué : {total_pages} pages, seules les {MAX_PDF_PAGES} premières traitées")
        result += f"\n\n[... tronqué à {MAX_PDF_PAGES} pages]"

    return result


def _extract_docx(file_path: Path) -> str:
    """Extract text from Word documents using python-docx."""
    try:
        from docx import Document
    except ImportError:
        logger.warning("python-docx not installed, cannot extract DOCX")
        return "[DOCX extraction unavailable - install python-docx]"

    doc = Document(file_path)
    text_parts = []

    for para in doc.paragraphs:
        if para.text.strip():
            text_parts.append(para.text)

    return "\n\n".join(text_parts)


def _extract_pptx(file_path: Path) -> str:
    """Extrait le texte d'une présentation, diapositive par diapositive.

    `python-pptx` est déjà une dépendance du projet (génération de documents) :
    accepter un .pptx à l'indexation sans jamais l'extraire n'avait donc aucune
    raison technique. Les notes du présentateur sont incluses, elles portent
    souvent l'essentiel du propos.
    """
    from pptx import Presentation

    presentation = Presentation(str(file_path))
    morceaux: list[str] = []

    for numero, diapositive in enumerate(presentation.slides, start=1):
        lignes = [
            forme.text_frame.text.strip()
            for forme in diapositive.shapes
            if forme.has_text_frame and forme.text_frame.text.strip()
        ]
        notes = ""
        if diapositive.has_notes_slide:
            notes = diapositive.notes_slide.notes_text_frame.text.strip()
        if not lignes and not notes:
            continue
        bloc = f"--- Diapositive {numero} ---"
        if lignes:
            bloc += "\n" + "\n".join(lignes)
        if notes:
            bloc += f"\n[Notes] {notes}"
        morceaux.append(bloc)

    return "\n\n".join(morceaux)


def _extract_xlsx(file_path: Path) -> str:
    """Extrait le contenu d'un fichier Excel en format texte tabulaire."""
    try:
        from openpyxl import load_workbook
    except ImportError:
        logger.warning("openpyxl not installed, cannot extract XLSX")
        return "[XLSX extraction unavailable - install openpyxl]"

    wb = load_workbook(file_path, read_only=True, data_only=True)
    output: list[str] = []
    lignes_lues = 0
    lignes_ignorees = 0

    for sheet in wb.sheetnames:
        ws = wb[sheet]
        output.append(f"## Feuille : {sheet}\n")
        for row in ws.iter_rows(values_only=True):
            if lignes_lues >= MAX_XLSX_LIGNES:
                lignes_ignorees += 1
                continue
            cells = [str(c) if c is not None else "" for c in row]
            output.append(" | ".join(cells))
            lignes_lues += 1
    wb.close()

    if lignes_ignorees:
        # Annoncé dans le texte lui-même, comme pour le PDF et le CSV : le
        # modèle doit pouvoir dire à l'utilisateur qu'il n'a pas tout vu.
        output.append(
            f"\n[Tableur tronqué : {lignes_lues} lignes transmises sur "
            f"{lignes_lues + lignes_ignorees}. Les suivantes n'ont pas été lues.]"
        )
        logger.info(
            "XLSX tronqué : %d lignes transmises, %d ignorées",
            lignes_lues, lignes_ignorees,
        )

    return "\n".join(output)


def chunk_text(
    text: str,
    chunk_size: int = 1000,
    overlap: int = 200,
    separator: str = "\n\n",
) -> Generator[str, None, None]:
    """
    Split text into overlapping chunks for embedding.

    Args:
        text: Text to chunk
        chunk_size: Target size of each chunk (characters)
        overlap: Overlap between chunks
        separator: Preferred split point

    Yields:
        Text chunks
    """
    if len(text) <= chunk_size:
        yield text
        return

    # Try to split on paragraphs first
    paragraphs = text.split(separator)

    current_chunk = ""
    for para in paragraphs:
        # If adding this paragraph exceeds chunk_size, yield current and start new
        if len(current_chunk) + len(para) + len(separator) > chunk_size:
            if current_chunk:
                yield current_chunk.strip()

            # Start new chunk with overlap from previous
            if len(current_chunk) > overlap:
                # Take last `overlap` characters
                current_chunk = current_chunk[-overlap:] + separator + para
            else:
                current_chunk = para
        else:
            if current_chunk:
                current_chunk += separator + para
            else:
                current_chunk = para

    # Yield remaining
    if current_chunk.strip():
        yield current_chunk.strip()


def get_file_metadata(file_path: Path) -> dict:
    """
    Get metadata about a file.

    Args:
        file_path: Path to file

    Returns:
        Dict with file metadata
    """
    stat = file_path.stat()
    mime_type, _ = mimetypes.guess_type(str(file_path))

    return {
        "name": file_path.name,
        "extension": file_path.suffix.lower(),
        "size": stat.st_size,
        "mime_type": mime_type,
        "modified_at": stat.st_mtime,
    }
