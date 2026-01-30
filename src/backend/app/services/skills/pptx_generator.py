"""
THÉRÈSE v2 - PowerPoint Generator Skill

Génère des présentations PowerPoint (.pptx) avec le style Synoptïa.
"""

import logging
import re
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Cm, Inches, Pt

from app.services.skills.base import BaseSkill, FileFormat, SkillParams, SkillResult

logger = logging.getLogger(__name__)


# Palette Synoptïa
SYNOPTIA_COLORS = {
    "background": RGBColor(0x0B, 0x12, 0x26),
    "surface": RGBColor(0x13, 0x1B, 0x35),
    "text": RGBColor(0xE6, 0xED, 0xF7),
    "muted": RGBColor(0xA9, 0xB8, 0xD8),
    "primary": RGBColor(0x17, 0x33, 0xA6),
    "accent_cyan": RGBColor(0x22, 0xD3, 0xEE),
    "accent_magenta": RGBColor(0xE1, 0x1D, 0x8D),
}


class PptxSkill(BaseSkill):
    """
    Skill de génération de présentations PowerPoint.

    Crée des présentations .pptx professionnelles avec le style Synoptïa dark.
    """

    skill_id = "pptx-pro"
    name = "Présentation PowerPoint"
    description = "Génère une présentation PowerPoint avec le style Synoptïa"
    output_format = FileFormat.PPTX

    def __init__(self, output_dir: Path):
        super().__init__(output_dir)

    async def execute(self, params: SkillParams) -> SkillResult:
        """
        Génère une présentation PowerPoint à partir du contenu.

        Args:
            params: Paramètres incluant titre et contenu

        Returns:
            Résultat avec chemin vers le fichier généré
        """
        file_id = self.generate_file_id()
        output_path = self.get_output_path(file_id, params.title)

        # Créer la présentation
        prs = Presentation()
        prs.slide_width = Inches(13.333)  # 16:9
        prs.slide_height = Inches(7.5)

        # Parser le contenu et créer les slides
        slides_content = self._parse_content(params.content)

        # Slide de titre
        self._add_title_slide(prs, params.title)

        # Slides de contenu
        for slide_data in slides_content:
            self._add_content_slide(prs, slide_data)

        # Slide de fin
        self._add_end_slide(prs)

        # Sauvegarder
        prs.save(str(output_path))

        # Calculer la taille
        file_size = output_path.stat().st_size

        logger.info(f"Generated PPTX: {output_path} ({file_size} bytes)")

        return SkillResult(
            file_id=file_id,
            file_path=output_path,
            file_name=output_path.name,
            file_size=file_size,
            mime_type=self.get_mime_type(),
            format=self.output_format,
        )

    def get_system_prompt_addition(self) -> str:
        """Instructions pour le LLM pour générer du contenu PowerPoint."""
        return """
## Instructions pour génération de présentation PowerPoint

Génère le contenu de la présentation en format structuré :

Pour chaque slide, utilise ce format :
---
# Titre de la slide
- Point 1
- Point 2
- Point 3
---

Règles :
1. **Titre** : Ligne commençant par # (max 8 mots)
2. **Points** : 3 à 5 points par slide (format liste -)
3. **Concision** : Max 10 mots par point
4. **Structure recommandée** :
   - Slide 1 : Contexte / Problématique
   - Slides 2-4 : Points clés / Solutions
   - Slide finale : Conclusion / Call-to-action

Style : impactant, visuel, orienté action.
Évite les phrases longues - privilégie les mots-clés.
"""

    def _parse_content(self, content: str) -> list[dict[str, Any]]:
        """
        Parse le contenu en structure de slides.

        Args:
            content: Contenu généré par le LLM

        Returns:
            Liste de dictionnaires avec titre et points pour chaque slide
        """
        slides = []
        current_slide = None

        # Découper par délimiteurs --- ou par titres #
        blocks = re.split(r'\n---\n|\n-{3,}\n', content)

        for block in blocks:
            block = block.strip()
            if not block:
                continue

            lines = block.split('\n')
            title = None
            points = []

            for line in lines:
                line = line.strip()
                if not line:
                    continue

                # Titre de slide
                if line.startswith('# '):
                    title = line[2:].strip()
                elif line.startswith('## '):
                    title = line[3:].strip()
                # Points
                elif line.startswith('- ') or line.startswith('* '):
                    points.append(line[2:].strip())
                elif re.match(r'^\d+\.\s', line):
                    points.append(re.sub(r'^\d+\.\s', '', line).strip())
                # Texte normal (traiter comme point si pas de titre encore)
                elif title and line:
                    points.append(line)

            if title or points:
                slides.append({
                    "title": title or "Slide",
                    "points": points[:6],  # Max 6 points par slide
                })

        return slides

    def _add_title_slide(self, prs: Presentation, title: str) -> None:
        """Ajoute la slide de titre."""
        blank_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(blank_layout)

        # Fond sombre
        self._set_background(slide, SYNOPTIA_COLORS["background"])

        # Titre centré
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5)
        )
        title_frame = title_box.text_frame
        title_frame.paragraphs[0].text = title
        title_frame.paragraphs[0].font.name = "Outfit"
        title_frame.paragraphs[0].font.size = Pt(54)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = SYNOPTIA_COLORS["text"]
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Sous-titre
        subtitle_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(4.2), Inches(12.333), Inches(0.5)
        )
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.paragraphs[0].text = "Synoptïa - L'entrepreneur augmenté"
        subtitle_frame.paragraphs[0].font.name = "Inter"
        subtitle_frame.paragraphs[0].font.size = Pt(18)
        subtitle_frame.paragraphs[0].font.italic = True
        subtitle_frame.paragraphs[0].font.color.rgb = SYNOPTIA_COLORS["accent_cyan"]
        subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    def _add_content_slide(self, prs: Presentation, slide_data: dict[str, Any]) -> None:
        """Ajoute une slide de contenu."""
        blank_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(blank_layout)

        # Fond sombre
        self._set_background(slide, SYNOPTIA_COLORS["background"])

        # Titre
        title_box = slide.shapes.add_textbox(
            Inches(0.75), Inches(0.5), Inches(11.833), Inches(1)
        )
        title_frame = title_box.text_frame
        title_frame.paragraphs[0].text = slide_data["title"]
        title_frame.paragraphs[0].font.name = "Outfit"
        title_frame.paragraphs[0].font.size = Pt(36)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = SYNOPTIA_COLORS["text"]

        # Barre accent sous le titre
        accent_bar = slide.shapes.add_shape(
            1,  # Rectangle
            Inches(0.75), Inches(1.4), Inches(2), Inches(0.05)
        )
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = SYNOPTIA_COLORS["accent_cyan"]
        accent_bar.line.fill.background()

        # Points
        content_box = slide.shapes.add_textbox(
            Inches(0.75), Inches(1.8), Inches(11.833), Inches(5)
        )
        content_frame = content_box.text_frame
        content_frame.word_wrap = True

        for i, point in enumerate(slide_data.get("points", [])):
            if i == 0:
                p = content_frame.paragraphs[0]
            else:
                p = content_frame.add_paragraph()

            # Bullet avec accent
            p.text = f"  {point}"
            p.font.name = "Inter"
            p.font.size = Pt(24)
            p.font.color.rgb = SYNOPTIA_COLORS["text"]
            p.space_before = Pt(16)
            p.level = 0

    def _add_end_slide(self, prs: Presentation) -> None:
        """Ajoute la slide de fin."""
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)

        # Fond sombre
        self._set_background(slide, SYNOPTIA_COLORS["background"])

        # Texte de fin
        text_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(3), Inches(12.333), Inches(1.5)
        )
        text_frame = text_box.text_frame

        text_frame.paragraphs[0].text = "Merci"
        text_frame.paragraphs[0].font.name = "Outfit"
        text_frame.paragraphs[0].font.size = Pt(48)
        text_frame.paragraphs[0].font.bold = True
        text_frame.paragraphs[0].font.color.rgb = SYNOPTIA_COLORS["accent_cyan"]
        text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Sous-texte
        p = text_frame.add_paragraph()
        p.text = "Généré par THÉRÈSE - Synoptïa"
        p.font.name = "Inter"
        p.font.size = Pt(14)
        p.font.color.rgb = SYNOPTIA_COLORS["muted"]
        p.alignment = PP_ALIGN.CENTER

    def _set_background(self, slide, color: RGBColor) -> None:
        """Définit la couleur de fond d'une slide."""
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color
