"""
THERESE v2 - Code Executor pour Skills Office

Exécution sandboxée de code Python généré par le LLM.
Approche code-execution : LLM -> code Python -> exécution -> fichier.
"""

import ast
import asyncio
import logging
import multiprocessing
import queue
import re
from abc import abstractmethod
from pathlib import Path
from typing import Any

from app.services.skills.base import BaseSkill, SkillParams, SkillResult

logger = logging.getLogger(__name__)

# Timeout d'exécution en secondes
EXECUTION_TIMEOUT = 30

# Palette Synoptia pour injection dans le namespace
SYNOPTIA_COLORS = {
    "background": "#0B1226",
    "surface": "#131B35",
    "text": "#E6EDF7",
    "muted": "#A9B8D8",
    "primary": "#2451FF",
    "accent_cyan": "#22D3EE",
    "accent_magenta": "#E11D8D",
    "header_bg": "0F1E6D",
    "header_text": "E6EDF7",
    "row_alt": "F5F7FA",
    "input_blue": "3B82F6",
    "formula_black": "1A1A2E",
    "link_green": "22C55E",
    "heading": "0F1E6D",
    "body": "1A1A2E",
}

# Patterns bloqués dans le code généré
BLOCKED_PATTERNS: list[str] = [
    r"\bos\.",
    r"\bos\b\s*\(",
    r"\bsys\.",
    r"\bsys\b\s*\(",
    r"\bsubprocess\b",
    r"\bshutil\b",
    r"\bsocket\b",
    r"\brequests\b",
    r"\burllib\b",
    r"\b__import__\b",
    r"\beval\s*\(",
    r"\bexec\s*\(",
    r"\bcompile\s*\(",
    r"\bglobals\s*\(",
    r"\blocals\s*\(",
    r"\bgetattr\s*\(",
    r"\bsetattr\s*\(",
    r"\bdelattr\s*\(",
    r"\bbreakpoint\s*\(",
    r"\binput\s*\(",
    # US-001 : bloquer l'introspection par dunders, qui permet l'évasion
    # classique du namespace restreint, ex. ().__class__.__bases__[0].__subclasses__().
    # Le code de génération de document (openpyxl/docx/pptx) n'en a aucun usage.
    r"__\w+__",
]

# Imports autorisés par format
ALLOWED_IMPORTS: dict[str, set[str]] = {
    "xlsx": {
        "openpyxl",
        "openpyxl.Workbook",
        "openpyxl.styles",
        "openpyxl.styles.Font",
        "openpyxl.styles.PatternFill",
        "openpyxl.styles.Alignment",
        "openpyxl.styles.Border",
        "openpyxl.styles.Side",
        "openpyxl.styles.numbers",
        "openpyxl.chart",
        "openpyxl.chart.BarChart",
        "openpyxl.chart.LineChart",
        "openpyxl.chart.PieChart",
        "openpyxl.chart.Reference",
        "openpyxl.utils",
        "openpyxl.utils.get_column_letter",
        # pandas retiré (passe 4, frontière de confiance) : read_csv /
        # to_excel ne passent pas par open() et lisaient ~/.therese, ~/.ssh
        # ou une URL HTTP. Le prompt du skill ne le proposait pas.
        "datetime",
        "json",
        "re",
        "math",
        "decimal",
        "decimal.Decimal",
        "time",
        "random",
        "copy",
        "string",
        "textwrap",
        "itertools",
        "collections",
    },
    "docx": {
        "docx",
        "docx.Document",
        "docx.shared",
        "docx.shared.Cm",
        "docx.shared.Pt",
        "docx.shared.Inches",
        "docx.shared.RGBColor",
        "docx.enum.text",
        "docx.enum.text.WD_ALIGN_PARAGRAPH",
        "docx.enum.style",
        "docx.enum.style.WD_STYLE_TYPE",
        "docx.enum.table",
        "docx.enum.table.WD_TABLE_ALIGNMENT",
        "docx.enum.section",
        "docx.oxml.ns",
        "docx.oxml.ns.qn",
        "datetime",
        "json",
        "re",
        "math",
        "decimal",
        "decimal.Decimal",
        "time",
        "random",
        "copy",
        "string",
        "textwrap",
        "itertools",
        "collections",
    },
    "pptx": {
        "pptx",
        "pptx.Presentation",
        "pptx.util",
        "pptx.util.Inches",
        "pptx.util.Pt",
        "pptx.util.Cm",
        "pptx.util.Emu",
        "pptx.dml.color",
        "pptx.dml.color.RGBColor",
        "pptx.enum.text",
        "pptx.enum.text.PP_ALIGN",
        "pptx.enum.text.MSO_ANCHOR",
        "pptx.enum.shapes",
        "pptx.enum.chart",
        "datetime",
        "json",
        "re",
        "math",
        "decimal",
        "decimal.Decimal",
        "time",
        "random",
        "copy",
        "string",
        "textwrap",
        "itertools",
        "collections",
    },
}


class CodeExecutionError(Exception):
    """Erreur lors de l'exécution du code généré."""

    pass


class LivrableInexploitable(CodeExecutionError):
    """Le modèle n'a rien produit d'utilisable : on refuse de livrer un faux.

    Revue 30/08/2026 : le repli fabriquait un classeur A/B/C, collait le
    script Python sur une diapositive, ou encapsulait un refus dans une
    page HTML. Un échec franc (skill_file_error) vaut mieux.
    """


# Seuils minimum de contenu pour valider un document généré par code-execution.
# Comptés hors coquille : titre, en-têtes inventés, slide Merci, pied de page.
MIN_CONTENT_ELEMENTS = {
    "docx": 3,  # au moins 3 paragraphes non vides
    "pptx": 2,  # titre + au moins une slide de contenu (Merci ne compte pas)
    "xlsx": 1,  # au moins 1 ligne de données, hors titre / en-tête / pied
}


def _ligne_est_du_python(ligne: str) -> bool:
    """True si la ligne est une instruction Python de générateur Office."""
    s = ligne.strip()
    if not s:
        return False
    if re.match(r"^(from\s+\S+\s+import|import\s+[A-Za-z_.]+)", s):
        return True
    if re.match(r"^(def|class|async def)\s+\w+", s):
        return True
    if re.match(r"^(print|raise|return|pass|break|continue)\b", s):
        return True
    if re.search(
        r"\b(Workbook|load_workbook|Document|Presentation|MSO_AUTO_SHAPE_TYPE)\s*[\(\.]",
        s,
    ):
        return True
    return bool(
        re.search(
            r"\.(save|add_slide|add_heading|add_paragraph|merge_cells|add_shape)\s*\(",
            s,
        )
    )


def retirer_code_python(contenu: str) -> str:
    """Retire les blocs de code avant un repli Markdown.

    Incident 30/08 : un ```python jamais refermé laissait le script dans le
    Word / la diapositive (Soso finding 2). On conserve la prose Markdown
    qui précède ou suit le bloc (BUG-135 : une fence orpheline n'avale
    plus la fin du document).
    """
    texte = re.sub(
        r"```(?:python|py|javascript|js|bash|sh|json|html|css|xml|sql|yaml|yml)?[^\n]*\n.*?```",
        "\n",
        contenu,
        flags=re.DOTALL,
    )
    # Fence orpheline : stop au prochain titre Markdown, sinon jusqu'à la fin.
    texte = re.sub(
        r"```(?:python|py)[^\n]*\n.*?(?=\n#{1,3} |\Z)",
        "\n",
        texte,
        flags=re.DOTALL,
    )
    lignes_out: list[str] = []
    for ligne in texte.splitlines():
        s = ligne.strip()
        if s.startswith("```"):
            continue
        if _ligne_est_du_python(s):
            continue
        lignes_out.append(ligne)
    return "\n".join(lignes_out).strip()


def contenu_ressemble_a_du_python(texte: str) -> bool:
    """True si le texte restant est encore du script, pas du Markdown métier."""
    lignes = [ligne.strip() for ligne in texte.splitlines() if ligne.strip()]
    if not lignes:
        return False
    hits = sum(1 for ligne in lignes if _ligne_est_du_python(ligne))
    return hits >= 3 or (hits >= 1 and hits / len(lignes) >= 0.3)


def _contenu_repli_exploitable(texte: str) -> bool:
    """Le repli Markdown n'a le droit de tourner que s'il reste du fond.

    Un titre + un paragraphe (`# Doc\\n\\nContenu.`) est un livrable. Une
    phrase d'intro orpheline après un script tronqué (« Voici le document : »)
    n'en est pas un : pas de structure, trop court.
    """
    if not texte or not texte.strip():
        return False
    if contenu_ressemble_a_du_python(texte):
        return False
    return bool(
        re.search(r"(?m)^#{1,3} ", texte)
        or re.search(r"(?m)^\s*\|.+\|", texte)
        or re.search(r"(?m)^[-*]\s+\S", texte)
        or len(texte.strip()) >= 120
    )


def _pptx_textes_slides(prs: Any) -> list[str]:
    textes: list[str] = []
    for slide in prs.slides:
        morceaux: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                morceaux.append(shape.text_frame.text.strip())
        textes.append("\n".join(m for m in morceaux if m))
    return textes


def _premiere_ligne_slide(texte: str) -> str:
    if not texte.strip():
        return ""
    return texte.strip().split("\n", 1)[0].strip()


def _pptx_cycle_repete(textes: list[str]) -> bool:
    """True si les slides de contenu répètent un cycle (P1 P2 P3 P1 P2 P3…).

    Constat du 30/08 : une demande en 3 parties sortait en 10 diapositives
    où 1-2-3 apparaissaient trois fois (nb_slides forcé à 10, le modèle
    remplissait en répétant).
    """
    corps = list(textes)
    if corps and re.match(r"^merci\b", _premiere_ligne_slide(corps[-1]), re.I):
        corps = corps[:-1]
    if len(corps) > 1:
        corps = corps[1:]
    if len(corps) < 6:
        return False
    cles = [_premiere_ligne_slide(t) for t in corps]
    n = len(cles)
    for k in range(2, n // 2 + 1):
        tours = n // k
        if tours < 2:
            continue
        cycle = cles[:k]
        if len(set(cycle)) < 2:
            continue
        reconstruit = cycle * tours + cycle[: n % k]
        if reconstruit == cles:
            return True
    return False


def _xlsx_lignes_de_donnees(ws: Any) -> int:
    """Compte les lignes métier, pas le titre, l'en-tête ni le pied.

    La garde d'avant comptait « Tableau » + A/B/C = 2 et laissait passer
    la coquille (revue 30/08, grok-chemin-nominal finding 2).
    """
    lignes: list[list[Any]] = []
    for row in ws.iter_rows(max_row=200):
        valeurs = [c.value for c in row if c.value not in (None, "")]
        if not valeurs:
            continue
        texte = " ".join(str(v) for v in valeurs)
        if re.search(r"g[ée]n[ée]r[ée]\s+par", texte, re.I):
            continue
        lignes.append(valeurs)
    if not lignes:
        return 0
    indice_entetes = 1 if len(lignes[0]) == 1 else 0
    if indice_entetes >= len(lignes):
        return 0
    entetes = [str(v) for v in lignes[indice_entetes]]
    donnees = lignes[indice_entetes + 1 :]
    if entetes == ["A", "B", "C"] and not donnees:
        return 0
    return len(donnees)


def _validate_document_content(output_path: str, format_type: str) -> bool:
    """
    Vérifie qu'un document généré contient assez de contenu.

    Retourne True si le document est suffisamment riche, False sinon.
    Une erreur de lecture refuse le fichier (fail-closed, revue 30/08) :
    on ne livre pas un objet qu'on n'a pas pu relire.
    """
    min_elements = MIN_CONTENT_ELEMENTS.get(format_type, 1)
    path = Path(output_path)

    try:
        if format_type == "docx":
            from docx import Document
            doc = Document(str(path))
            non_empty = sum(
                1
                for p in doc.paragraphs
                if p.text.strip() and not _ligne_est_du_python(p.text)
            )
            logger.debug(
                "Validation DOCX : %d paragraphes non vides (min=%d)",
                non_empty, min_elements,
            )
            return non_empty >= min_elements

        elif format_type == "pptx":
            from pptx import Presentation
            prs = Presentation(str(path))
            textes = _pptx_textes_slides(prs)
            if _pptx_cycle_repete(textes):
                logger.warning(
                    "Validation PPTX : cycle de diapositives répété, livrable refusé"
                )
                return False
            hors_merci = [
                t
                for t in textes
                if t.strip()
                and not re.match(r"^merci\b", _premiere_ligne_slide(t), re.I)
            ]
            logger.debug(
                "Validation PPTX : %d slides hors Merci (min=%d)",
                len(hors_merci), min_elements,
            )
            return len(hors_merci) >= min_elements

        elif format_type == "xlsx":
            from openpyxl import load_workbook
            wb = load_workbook(str(path), read_only=True)
            try:
                total = 0
                for ws in wb.worksheets:
                    total += _xlsx_lignes_de_donnees(ws)
            finally:
                wb.close()
            logger.debug(
                "Validation XLSX : %d lignes de données (min=%d)",
                total, min_elements,
            )
            return total >= min_elements

        elif format_type == "html":
            texte = path.read_text(encoding="utf-8", errors="replace")
            return "<html" in texte.lower() and "</html>" in texte.lower()

    except Exception as e:
        logger.warning("Erreur validation contenu %s : %s", format_type, e)
        return False

    return False


def extract_python_code(llm_response: str) -> str | None:
    """
    Extrait le code Python d'un bloc ```python``` dans la réponse LLM.

    Gère aussi les réponses tronquées où le ``` fermant est absent
    (ex: réponses longues coupées par max_tokens).

    Args:
        llm_response: Réponse complète du LLM

    Returns:
        Code Python extrait ou None si pas de bloc trouvé
    """
    # 1. Chercher un bloc complet ```python ... ```
    pattern = r"```python\s*\n(.*?)```"
    match = re.search(pattern, llm_response, re.DOTALL)
    if match:
        code = match.group(1).strip()
        if code:
            return code

    # 2. Essayer aussi ```py ... ```
    pattern_py = r"```py\s*\n(.*?)```"
    match_py = re.search(pattern_py, llm_response, re.DOTALL)
    if match_py:
        code = match_py.group(1).strip()
        if code:
            return code

    # 3. Fallback : bloc ouvert sans ``` fermant (réponse tronquée)
    pattern_open = r"```python\s*\n(.*)"
    match_open = re.search(pattern_open, llm_response, re.DOTALL)
    if match_open:
        code = match_open.group(1).strip()
        if code:
            logger.warning(
                "Code Python extrait sans bloc fermant (réponse tronquée, %d chars)",
                len(code),
            )
            return code

    # 4. Idem pour ```py sans fermant
    pattern_py_open = r"```py\s*\n(.*)"
    match_py_open = re.search(pattern_py_open, llm_response, re.DOTALL)
    if match_py_open:
        code = match_py_open.group(1).strip()
        if code:
            logger.warning(
                "Code Python extrait sans bloc fermant (réponse tronquée, %d chars)",
                len(code),
            )
            return code

    # 5. Dernier recours : détecter du code Python brut (sans bloc markdown)
    # Certains modèles (ex: Gemini) renvoient du code sans ```python```
    # On cherche des marqueurs forts de code python-docx/pptx/openpyxl
    code_markers = [
        r'from\s+(?:docx|pptx|openpyxl)\s+import',
        r'(?:Document|Presentation|Workbook)\s*\(',
        r'\.add_slide\s*\(',
        r'\.add_heading\s*\(',
        r'\.save\s*\(',
    ]
    marker_count = sum(1 for m in code_markers if re.search(m, llm_response))
    if marker_count >= 2:
        # Extraire les lignes qui ressemblent à du code Python
        # (ignorer les lignes de texte libre avant/après le code)
        lines = llm_response.split('\n')
        code_lines: list[str] = []
        in_code = False
        for line in lines:
            stripped = line.strip()
            # Début du code : import ou assignation ou commentaire Python
            if not in_code:
                if (stripped.startswith(('from ', 'import ', '#'))
                    or re.match(r'^[a-zA-Z_]\w*\s*=\s*', stripped)):
                    in_code = True
                    code_lines.append(line)
            else:
                # Fin du code : ligne vide après du contenu, ou texte narratif
                if stripped and not stripped.startswith('#') and not any(c in stripped for c in '=()[]{}.:,+-*/_"\'\\') and len(stripped.split()) > 5:
                    # Ligne de texte narratif (pas de code) - on arrête si on a déjà du code
                    if len(code_lines) > 5:
                        break
                else:
                    code_lines.append(line)

        if code_lines:
            code = '\n'.join(code_lines).strip()
            try:
                ast.parse(code)
                logger.warning(
                    "Code Python détecté sans bloc markdown (%d lignes, %d markers)",
                    len(code_lines),
                    marker_count,
                )
                return code
            except SyntaxError:
                # Tenter réparation
                repaired = repair_truncated_code(code)
                if repaired:
                    logger.warning(
                        "Code Python sans bloc markdown réparé (%d lignes)",
                        len(code_lines),
                    )
                    return repaired

    return None


def repair_truncated_code(code: str) -> str | None:
    """
    Tente de réparer du code Python tronqué en retirant les lignes
    incomplètes à la fin jusqu'à obtenir un code syntaxiquement valide.

    Si la réparation supprime l'appel .save(output_path), celui-ci est
    rajouté automatiquement en détectant le nom de la variable du document
    (wb, doc, prs, document, workbook, presentation).

    Args:
        code: Code Python potentiellement tronqué

    Returns:
        Code réparé ou None si impossible
    """
    # D'abord vérifier si le code est déjà valide
    try:
        ast.parse(code)
        return _ensure_save_call(code)
    except SyntaxError:
        pass

    # Retirer les lignes une par une depuis la fin
    lines = code.split("\n")
    for i in range(len(lines), 0, -1):
        candidate = "\n".join(lines[:i]).rstrip()
        if not candidate:
            continue
        try:
            ast.parse(candidate)
            removed = len(lines) - i
            logger.info(
                "Code tronqué réparé : %d lignes retirées sur %d",
                removed,
                len(lines),
            )
            # Ajouter .save(output_path) si absent après réparation
            candidate = _ensure_save_call(candidate)
            return candidate
        except SyntaxError:
            continue

    return None


def _forcer_sauvegarde_vers_sortie(code: str) -> str:
    """Réécrit tout appel .save(...) en .save(output_path).

    Passe 4 : la réécriture regex ne couvrait qu'un littéral
    (`.save("facture.docx")`). `chemin = "/Users/ludo/..."; wb.save(chemin)`
    passait au travers, et `_ensure_save_call` s'arrêtait dès qu'un
    `.save(` existait. Un AST ne se laisse pas tromper par une variable.
    """
    try:
        tree = ast.parse(code)
    except SyntaxError:
        return code

    class _SaveVersSortie(ast.NodeTransformer):
        def visit_Call(self, node: ast.Call) -> ast.AST:
            self.generic_visit(node)
            if isinstance(node.func, ast.Attribute) and node.func.attr == "save":
                return ast.Call(
                    func=node.func,
                    args=[ast.Name(id="output_path", ctx=ast.Load())],
                    keywords=[],
                )
            return node

    nouveau = _SaveVersSortie().visit(tree)
    ast.fix_missing_locations(nouveau)
    return ast.unparse(nouveau)


def _ensure_save_call(code: str) -> str:
    """
    Vérifie que le code contient un appel .save(output_path).
    Si absent, détecte le nom de la variable du document et l'ajoute.

    Args:
        code: Code Python syntaxiquement valide

    Returns:
        Code avec .save(output_path) garanti
    """
    code = _forcer_sauvegarde_vers_sortie(code)
    if re.search(r"\.save\s*\(\s*output_path\s*\)", code):
        return code

    # Détecter le nom de la variable du document principal
    # Patterns courants : wb = Workbook(), doc = Document(), prs = Presentation()
    # Inclut aussi load_workbook() pour les fichiers Excel existants
    doc_var = None
    for pattern in [
        r"(\w+)\s*=\s*Workbook\s*\(",
        r"(\w+)\s*=\s*load_workbook\s*\(",
        r"(\w+)\s*=\s*Document\s*\(",
        r"(\w+)\s*=\s*Presentation\s*\(",
    ]:
        match = re.search(pattern, code)
        if match:
            doc_var = match.group(1)
            break

    if doc_var:
        save_line = f"\n{doc_var}.save(output_path)\n"
        logger.warning(
            "Appel .save(output_path) manquant, "

            "ajout automatique : %s.save(output_path)",
            doc_var,
        )
        return code + save_line

    return code


def validate_code(code: str) -> tuple[bool, str]:
    """
    Valide la sécurité du code Python généré.

    Vérifie :
    1. Syntaxe Python valide (via ast.parse)
    2. Pas de patterns dangereux (os, sys, subprocess, etc.)
    3. Pas d'imports non autorisés

    Args:
        code: Code Python à valider

    Returns:
        Tuple (est_valide, message_erreur)
    """
    # 1. Vérifier la syntaxe
    try:
        ast.parse(code)
    except SyntaxError as e:
        return False, f"Erreur de syntaxe Python : {e}"

    # 2. Vérifier les patterns bloqués
    for pattern in BLOCKED_PATTERNS:
        if re.search(pattern, code):
            return False, f"Pattern interdit détecté : {pattern}"

    # 3. Vérifier que open() n'est utilisé qu'avec output_path
    # On autorise open() uniquement via les appels .save() des bibliothèques
    open_calls = re.findall(r"\bopen\s*\(", code)
    if open_calls:
        # Vérifier que open() est utilisé uniquement avec output_path
        # Pattern autorisé : open(output_path, ...) ou open(str(output_path), ...)
        safe_open = re.findall(
            r"\bopen\s*\(\s*(?:str\s*\(\s*)?output_path", code
        )
        if len(open_calls) != len(safe_open):
            return False, "open() n'est autorisé qu'avec output_path"

    return True, ""


def _validate_imports(code: str, format_type: str) -> tuple[bool, str]:
    """
    Vérifie que les imports sont autorisés pour le format donné.

    Args:
        code: Code Python à valider
        format_type: Format du fichier (xlsx, docx, pptx)

    Returns:
        Tuple (est_valide, message_erreur)
    """
    allowed = ALLOWED_IMPORTS.get(format_type, set())

    try:
        tree = ast.parse(code)
    except SyntaxError:
        return False, "Erreur de syntaxe"

    for node in ast.walk(tree):
        if isinstance(node, ast.Import):
            for alias in node.names:
                module_root = alias.name.split(".")[0]
                if module_root not in {m.split(".")[0] for m in allowed}:
                    return False, f"Import interdit : {alias.name}"
        elif isinstance(node, ast.ImportFrom) and node.module:
            module_root = node.module.split(".")[0]
            if module_root not in {m.split(".")[0] for m in allowed}:
                return False, f"Import interdit : from {node.module}"

    return True, ""


def _build_namespace(
    output_path: str, title: str, format_type: str, nb_slides: int = 10
) -> dict[str, Any]:
    """
    Construit le namespace d'exécution sandboxé.

    Args:
        output_path: Chemin de sauvegarde du fichier
        title: Titre du document
        format_type: Format du fichier (xlsx, docx, pptx)

    Returns:
        Namespace dict avec les variables et imports autorisés
    """
    import datetime
    import json
    import math
    import re as re_module
    from decimal import Decimal

    namespace: dict[str, Any] = {
        # Variables injectées
        "output_path": output_path,
        "title": title,
        "nb_slides": nb_slides,
        "SYNOPTIA_COLORS": SYNOPTIA_COLORS.copy(),
        # Modules communs
        "datetime": datetime,
        "json": json,
        "re": re_module,
        "math": math,
        "Decimal": Decimal,
        # Builtins restreints
        "__builtins__": {
            "print": print,
            "len": len,
            "range": range,
            "enumerate": enumerate,
            "zip": zip,
            "map": map,
            "filter": filter,
            "sorted": sorted,
            "reversed": reversed,
            "list": list,
            "dict": dict,
            "set": set,
            "tuple": tuple,
            "str": str,
            "int": int,
            "float": float,
            "bool": bool,
            "max": max,
            "min": min,
            "sum": sum,
            "abs": abs,
            "round": round,
            "isinstance": isinstance,
            "type": type,
            "hasattr": hasattr,
            "None": None,
            "True": True,
            "False": False,
            "ValueError": ValueError,
            "TypeError": TypeError,
            "KeyError": KeyError,
            "IndexError": IndexError,
            "Exception": Exception,
            "__import__": _restricted_import(format_type),
        },
    }

    # Imports spécifiques selon le format
    if format_type == "xlsx":
        import openpyxl
        from openpyxl import Workbook
        from openpyxl.chart import BarChart, LineChart, PieChart, Reference
        from openpyxl.styles import (
            Alignment,
            Border,
            Font,
            PatternFill,
            Side,
        )
        from openpyxl.utils import get_column_letter

        namespace.update(
            {
                "openpyxl": openpyxl,
                "Workbook": Workbook,
                "Font": Font,
                "PatternFill": PatternFill,
                "Alignment": Alignment,
                "Border": Border,
                "Side": Side,
                "BarChart": BarChart,
                "LineChart": LineChart,
                "PieChart": PieChart,
                "Reference": Reference,
                "get_column_letter": get_column_letter,
            }
        )
    elif format_type == "docx":
        import docx
        from docx import Document
        from docx.enum.text import WD_ALIGN_PARAGRAPH
        from docx.shared import Cm, Inches, Pt, RGBColor

        namespace.update(
            {
                "docx": docx,
                "Document": Document,
                "Cm": Cm,
                "Pt": Pt,
                "Inches": Inches,
                "RGBColor": RGBColor,
                "WD_ALIGN_PARAGRAPH": WD_ALIGN_PARAGRAPH,
            }
        )
    elif format_type == "pptx":
        import pptx
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
        from pptx.util import Cm, Emu, Inches, Pt

        namespace.update(
            {
                "pptx": pptx,
                "Presentation": Presentation,
                "Inches": Inches,
                "Pt": Pt,
                "Cm": Cm,
                "Emu": Emu,
                "RGBColor": RGBColor,
                "PP_ALIGN": PP_ALIGN,
                "MSO_ANCHOR": MSO_ANCHOR,
            }
        )

    return namespace


def _restricted_import(format_type: str):
    """
    Crée une fonction __import__ restreinte aux modules autorisés.

    Args:
        format_type: Format du fichier (xlsx, docx, pptx)

    Returns:
        Fonction __import__ sécurisée
    """
    import builtins

    allowed_roots = {m.split(".")[0] for m in ALLOWED_IMPORTS.get(format_type, set())}

    def safe_import(name, *args, **kwargs):
        root = name.split(".")[0]
        if root not in allowed_roots:
            raise ImportError(
                f"Import interdit : '{name}'. "
                f"Modules autorisés : {sorted(allowed_roots)}"
            )
        return builtins.__import__(name, *args, **kwargs)

    return safe_import


def _installer_garde_fs(output_path: str) -> Any:
    """Borne builtins.open au dossier de sortie, dans CE process.

    Passe 4 : `open()` du namespace sandboxé était absent (NameError),
    mais openpyxl / python-docx / zipfile utilisent le vrai
    `builtins.open` du process. Le sous-processus tourne sous
    l'utilisateur de Ludo : une lecture de ~/.ssh passait.

    On patche le process entier, pas le namespace : c'est le seul
    endroit qui intercepte les bibliothèques déjà importées. Renvoie
    un restaurateur : les tests appellent le worker dans le process
    pytest, il ne faut pas y laisser open() borné.
    """
    import builtins
    import io
    import os
    import tempfile

    reel_open = builtins.open
    reel_io_open = io.open
    dossier = Path(output_path).resolve().parent
    # openpyxl écrit un NamedTemporaryFile puis le relit. Sans ça, un
    # save légitime échoue : le temp atterrit dans /var/folders. On
    # redirige le dossier temp vers la sortie AVANT d'ouvrir quoi que
    # ce soit : écrire dans /tmp n'est plus un exutoire.
    ancien_tmpdir = os.environ.get("TMPDIR")
    ancien_tempdir = tempfile.tempdir
    os.environ["TMPDIR"] = str(dossier)
    tempfile.tempdir = str(dossier)

    def open_borne(file: Any, mode: str = "r", *args: Any, **kwargs: Any) -> Any:
        if isinstance(file, int):
            return reel_open(file, mode, *args, **kwargs)
        try:
            chemin = Path(file).resolve()
            chemin.relative_to(dossier)
        except (TypeError, ValueError, OSError) as exc:
            raise PermissionError(
                f"Accès hors du dossier de sortie interdit : {file}"
            ) from exc
        return reel_open(file, mode, *args, **kwargs)

    # ZipFile utilise io.open, lié à l'import, pas builtins.open.
    builtins.open = open_borne
    io.open = open_borne

    def restaurer() -> None:
        builtins.open = reel_open
        io.open = reel_io_open
        if ancien_tmpdir is None:
            os.environ.pop("TMPDIR", None)
        else:
            os.environ["TMPDIR"] = ancien_tmpdir
        tempfile.tempdir = ancien_tempdir

    return restaurer


def _installer_garde_reseau() -> Any:
    """Refuse toute ouverture de socket dans ce process.

    pandas.read_csv('https://…') parlait HTTP via urllib, pas open().
    Interdire le socket coupe ce canal, quelle que soit la bibliothèque.
    Renvoie un restaurateur (même raison que la garde FS).
    """
    import socket

    reel_socket = socket.socket

    def _refuser(*_args: Any, **_kwargs: Any) -> Any:
        raise PermissionError(
            "Réseau interdit dans le bac à sable des documents"
        )

    setattr(socket, "socket", _refuser)  # noqa: B010

    def restaurer() -> None:
        setattr(socket, "socket", reel_socket)  # noqa: B010

    return restaurer


# B-250 : ce qui reste dans l'environnement du bac à sable. Liste BLANCHE et
# non liste noire : une variable ajoutée demain au backend (une clé de
# fournisseur de plus, un jeton de session) est écartée sans qu'on ait à y
# penser, ce qui est précisément la défense en profondeur promise.
_ENV_CONSERVE = frozenset({
    # De quoi trouver l'interpréteur et ses bibliothèques
    "PATH",
    "PYTHONPATH",
    "PYTHONHOME",
    "PYTHONIOENCODING",
    "PYTHONUTF8",
    "VIRTUAL_ENV",
    # Dossiers de travail (la garde FS réécrit TMPDIR juste après)
    "HOME",
    "TMPDIR",
    "TEMP",
    "TMP",
    # Locale et fuseau : sans eux les dates et l'encodage changent de forme
    "LANG",
    "LANGUAGE",
    "TZ",
    # Windows : sans ces variables, os.path et les bibliothèques standard
    # perdent pied
    "SYSTEMROOT",
    "SYSTEMDRIVE",
    "WINDIR",
    "COMSPEC",
    "PATHEXT",
    "USERPROFILE",
    "APPDATA",
    "LOCALAPPDATA",
    "PROGRAMFILES",
    "PROGRAMDATA",
    "NUMBER_OF_PROCESSORS",
    "PROCESSOR_ARCHITECTURE",
    "OS",
})
_ENV_PREFIXES_CONSERVES = ("LC_",)


def _epurer_environnement() -> Any:
    """Retire de `os.environ` tout ce qui n'est pas nécessaire à la génération.

    B-250 : `spawn` donne à l'enfant un interpréteur neuf, donc rien de la
    MÉMOIRE du backend - mais l'ENVIRONNEMENT, lui, était hérité en entier.
    `encryption.get_db_key_hex()` lit `THERESE_DB_KEY` dans `os.environ`, et
    le backend s'y replie aussi pour les clés de fournisseurs posées par
    l'utilisateur (ANTHROPIC_API_KEY, OPENAI_API_KEY, GEMINI_API_KEY,
    GROQ_API_KEY, BRAVE_API_KEY, FAL_API_KEY, OPENROUTER_API_KEY). La promesse
    « une évasion ne donne pas accès aux secrets du process principal » ne
    tenait donc qu'à moitié.

    `multiprocessing` n'accepte pas d'environnement pour l'enfant : l'épuration
    se fait dans l'enfant, première instruction de la cible.

    Renvoie un restaurateur, pour la même raison que les gardes FS et réseau :
    les tests appellent le worker dans le process pytest, il ne faut pas y
    laisser un environnement amputé.
    """
    import os

    sauvegarde = dict(os.environ)
    for nom in list(os.environ):
        if nom in _ENV_CONSERVE or nom.startswith(_ENV_PREFIXES_CONSERVES):
            continue
        os.environ.pop(nom, None)

    def restaurer() -> None:
        os.environ.clear()
        os.environ.update(sauvegarde)

    return restaurer


def _run_generation_in_subprocess(
    code: str,
    output_path: str,
    title: str,
    format_type: str,
    nb_slides: int,
    result_queue: Any,
) -> None:
    """Cible du sous-process spawn (US-001).

    Exécutée dans un interpréteur frais qui n'hérite pas de la mémoire du
    backend (token de session, clé Fernet) : une évasion éventuelle du
    namespace restreint ne donne donc pas accès aux secrets du process
    principal. Le résultat (ok / error) est remonté via la queue.

    Passe 4 : on importe les bibliothèques Office AVEC le vrai open
    (sinon site-packages est hors du dossier de sortie), puis on
    installe les gardes, puis on réécrit les `.save(...)` vers
    output_path. Un code généré ne lit plus hors de ce qu'on lui
    donne, et n'écrit plus hors du dossier de sortie.
    """
    restaurer_reseau = None
    restaurer_fs = None
    restaurer_env = None
    try:
        restaurer_env = _epurer_environnement()
        restaurer_reseau = _installer_garde_reseau()
        namespace = _build_namespace(output_path, title, format_type, nb_slides)
        restaurer_fs = _installer_garde_fs(output_path)
        code = _forcer_sauvegarde_vers_sortie(code)
        compiled = compile(code, "<llm_generated>", "exec")
        exec(compiled, namespace)  # noqa: S102
        result_queue.put(("ok", ""))
    except Exception as e:  # noqa: BLE001 - tout est remonté au process parent
        result_queue.put(("error", f"{type(e).__name__}: {e}"))
    finally:
        if restaurer_fs is not None:
            restaurer_fs()
        if restaurer_reseau is not None:
            restaurer_reseau()
        if restaurer_env is not None:
            restaurer_env()


async def execute_sandboxed(
    code: str,
    output_path: str,
    title: str,
    format_type: str,
    nb_slides: int = 10,
) -> None:
    """
    Exécute du code Python dans un sous-process isolé avec timeout.

    Args:
        code: Code Python à exécuter
        output_path: Chemin de sauvegarde du fichier
        title: Titre du document
        format_type: Format du fichier (xlsx, docx, pptx)

    Raises:
        CodeExecutionError: Si l'exécution échoue
    """
    # 1. Valider la sécurité du code
    is_valid, error_msg = validate_code(code)
    if not is_valid:
        raise CodeExecutionError(f"Validation échouée : {error_msg}")

    # 2. Valider les imports
    is_valid_imports, import_error = _validate_imports(code, format_type)
    if not is_valid_imports:
        raise CodeExecutionError(f"Validation imports échouée : {import_error}")

    # 3. Exécuter dans un sous-process spawn isolé (US-001).
    # spawn => interpréteur neuf sans la mémoire du backend ; combiné au blocage
    # des dunders, une évasion ne peut pas atteindre les secrets du process.
    ctx = multiprocessing.get_context("spawn")
    result_queue: Any = ctx.Queue()
    process = ctx.Process(
        target=_run_generation_in_subprocess,
        args=(code, output_path, title, format_type, nb_slides, result_queue),
        daemon=True,
    )

    def _run_and_wait() -> tuple[str, str]:
        process.start()
        try:
            status, detail = result_queue.get(timeout=EXECUTION_TIMEOUT)
        except queue.Empty:
            status, detail = "timeout", ""
        finally:
            if process.is_alive():
                process.terminate()
            process.join(5)
        return status, detail

    status, detail = await asyncio.to_thread(_run_and_wait)

    if status == "timeout":
        raise CodeExecutionError(
            f"Timeout : l'exécution a dépassé {EXECUTION_TIMEOUT}s"
        )
    if status == "error":
        raise CodeExecutionError(f"Erreur d'exécution : {detail}")


class CodeGenSkill(BaseSkill):
    """
    Classe abstraite pour les skills qui génèrent du code Python.

    Étend BaseSkill avec une approche code-execution :
    1. Le LLM génère du code Python dans un bloc ```python```
    2. Le code est validé et exécuté dans une sandbox
    3. Si échec ou pas de code, fallback vers l'ancien parser
    """

    async def execute(self, params: SkillParams) -> SkillResult:
        """
        Exécute le skill avec approche code-execution.

        1. Extraire code Python de params.content
        2. Si code trouvé -> exécuter dans sandbox
        3. Si échec ou pas de code -> _fallback_execute()
        4. Vérifier que le fichier a été créé (taille > 0)

        Args:
            params: Paramètres de génération

        Returns:
            Résultat avec chemin vers le fichier généré
        """
        file_id = self.generate_file_id()
        output_path = self.get_output_path(file_id, params.title)

        # Tenter l'extraction du code Python
        code = extract_python_code(params.content)

        if code:
            # Si le code a une erreur de syntaxe, tenter une réparation
            try:
                ast.parse(code)
            except SyntaxError:
                logger.warning(
                    f"[{self.skill_id}] Code extrait avec erreur de syntaxe, "
                    f"tentative de réparation..."
                )
                repaired = repair_truncated_code(code)
                if repaired:
                    code = repaired
                else:
                    logger.warning(
                        f"[{self.skill_id}] Réparation impossible, "
                        f"fallback vers parser legacy"
                    )
                    code = None

        if code:
            # Patcher les appels .save() pour utiliser output_path
            # Le LLM écrit souvent doc.save("fichier.docx") au lieu de doc.save(output_path)
            code = re.sub(
                r'\.save\s*\(\s*["\'].*?["\']\s*\)',
                '.save(output_path)',
                code,
            )
            # Aussi patcher wb.save(...) pour xlsx
            code = re.sub(
                r'\.save\s*\(\s*(?:f["\'].*?["\']|["\'].*?["\'])\s*\)',
                '.save(output_path)',
                code,
            )

            # Ajouter .save(output_path) si absent (code tronqué par max_tokens
            # mais syntaxiquement valide, donc pas réparé par repair_truncated_code)
            code = _ensure_save_call(code)

            try:
                logger.info(
                    f"[{self.skill_id}] Code Python détecté, exécution sandboxée..."
                )
                await execute_sandboxed(
                    code=code,
                    output_path=str(output_path),
                    title=params.title,
                    format_type=self.output_format.value,
                    nb_slides=params.metadata.get("nb_slides", 10),
                )

                # Vérifier que le fichier a été créé et contient du contenu
                if output_path.exists() and output_path.stat().st_size > 0:
                    # BUG-043 : valider que le document n'est pas quasi vide
                    if not _validate_document_content(
                        str(output_path), self.output_format.value
                    ):
                        logger.warning(
                            f"[{self.skill_id}] Code exécuté mais document quasi vide "
                            f"(contenu insuffisant), fallback vers parser legacy"
                        )
                    else:
                        file_size = output_path.stat().st_size
                        logger.info(
                            f"[{self.skill_id}] Code-execution réussi : "
                            f"{output_path} ({file_size} bytes)"
                        )
                        return SkillResult(
                            file_id=file_id,
                            file_path=output_path,
                            file_name=output_path.name,
                            file_size=file_size,
                            mime_type=self.get_mime_type(),
                            format=self.output_format,
                        )
                else:
                    logger.warning(
                        f"[{self.skill_id}] Code exécuté mais fichier vide ou inexistant, "
                        f"fallback vers parser legacy"
                    )
            except CodeExecutionError as e:
                logger.warning(
                    f"[{self.skill_id}] Échec code-execution : {e}, "
                    f"fallback vers parser legacy"
                )
            except Exception as e:
                logger.warning(
                    f"[{self.skill_id}] Erreur inattendue code-execution : {e}, "
                    f"fallback vers parser legacy"
                )
        else:
            logger.info(
                f"[{self.skill_id}] Pas de bloc Python détecté, "
                f"fallback vers parser legacy"
            )

        # Repli Markdown seulement s'il reste du fond métier. Garder le
        # script (fence orpheline, revue 30/08) produisait un livrable faux
        # au lieu d'un échec franc.
        contenu_repli = retirer_code_python(params.content)
        if not _contenu_repli_exploitable(contenu_repli):
            raise LivrableInexploitable(
                "le contenu produit n'est pas exploitable"
            )

        logger.info(
            f"[{self.skill_id}] fallback vers parser legacy"
        )
        params = SkillParams(
            title=params.title,
            content=contenu_repli,
            template=params.template,
            metadata=params.metadata,
        )
        result = await self._fallback_execute(params, file_id, output_path)
        if (
            not result.file_path.exists()
            or result.file_path.stat().st_size == 0
            or not _validate_document_content(
                str(result.file_path), self.output_format.value
            )
        ):
            raise LivrableInexploitable(
                "le document de repli n'est pas exploitable"
            )
        return result

    def get_markdown_prompt_addition(self) -> str:
        """
        Instructions alternatives pour les modèles incapables de générer du code Python.

        Demande du Markdown structuré au lieu de code python-docx/pptx/openpyxl.
        Peut être surchargé par chaque generator pour des instructions spécifiques.
        """
        return """
Génère le contenu en Markdown bien structuré.
Utilise : # Titre, ## Sections, ### Sous-sections, listes, **gras**, *italique*, tableaux Markdown.
NE génère PAS de code Python. Écris directement le contenu textuel.
"""

    @abstractmethod
    async def _fallback_execute(
        self, params: SkillParams, file_id: str, output_path: Path
    ) -> SkillResult:
        """
        Fallback vers l'ancien parser si le code-execution échoue.

        Args:
            params: Paramètres de génération
            file_id: ID du fichier pré-généré
            output_path: Chemin de sortie pré-calculé

        Returns:
            Résultat avec chemin vers le fichier généré
        """
        pass
