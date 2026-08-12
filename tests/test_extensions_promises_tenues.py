"""Aucune extension acceptée à l'indexation ne doit être illisible.

Relevé par l'inventaire des capacités du 13/08/2026 : le contrôle d'entrée
accepte 52 extensions, le parseur n'en extrait que 44. Les 13 restantes
produisent le pire des retours possibles — l'indexation « réussit », un
document apparaît dans la liste, et rien n'a été lu. L'utilisateur n'a aucun
moyen de le savoir : ni erreur, ni avertissement, ni contenu vide visible.

Six d'entre elles sont du texte brut que le parseur saurait lire s'il en avait
été informé (`.cfg`, `.conf`, `.ini`, `.org`, `.scala`, `.tex`), une est un
format bureautique dont la bibliothèque est déjà une dépendance (`.pptx`), et
les six dernières sont des binaires que rien ne sait ouvrir.

Ce test est un gate durable : il interdit qu'une extension soit ajoutée à la
liste blanche sans que le parseur sache l'extraire. Accepter un fichier, c'est
promettre de le lire.
"""


def _extensions_extractibles() -> set[str]:
    """Ce que `extract_text` sait réellement traiter, d'après ses branches."""
    from app.services import file_parser as fp

    return (
        fp.TEXT_EXTENSIONS
        | fp.CODE_EXTENSIONS
        | fp.CSV_EXTENSIONS
        | fp.PRESENTATION_EXTENSIONS
        | {".pdf", ".docx", ".doc", ".xlsx"}
    )


class TestTouteExtensionAccepteeEstLisible:
    def test_aucune_extension_acceptee_n_est_illisible(self):
        from app.services.path_security import INDEXABLE_EXTENSIONS

        orphelines = INDEXABLE_EXTENSIONS - _extensions_extractibles()

        assert not orphelines, (
            "ces extensions sont acceptées à l'indexation mais aucun chemin du "
            "parseur ne sait les extraire : l'utilisateur croit avoir indexé un "
            f"document dont rien n'a été lu — {sorted(orphelines)}"
        )

    def test_les_formats_texte_courants_sont_reconnus(self):
        """Refuser un fichier de configuration ou un .tex serait absurde.

        Ce sont des fichiers texte : les lire ne demande aucune bibliothèque,
        seulement de les déclarer.
        """
        from app.services import file_parser as fp

        attendus = {".cfg", ".conf", ".ini", ".org", ".scala", ".tex"}
        connus = fp.TEXT_EXTENSIONS | fp.CODE_EXTENSIONS

        assert attendus <= connus, (
            f"formats texte encore ignorés : {sorted(attendus - connus)}"
        )

    def test_un_format_non_extractible_est_refuse_a_l_entree(self, tmp_path):
        """Mieux vaut un refus net qu'un succès mensonger.

        Les binaires que rien ne sait ouvrir doivent être rejetés au contrôle
        d'entrée, avec un motif lisible, plutôt qu'indexés à vide.
        """
        from app.services.path_security import INDEXABLE_EXTENSIONS

        for extension in (".odt", ".ods", ".odp", ".ppt", ".rtf", ".xls"):
            assert extension not in INDEXABLE_EXTENSIONS, (
                f"{extension} est encore accepté alors que rien ne sait "
                "l'extraire : l'indexation produira un document vide sans le dire"
            )


class TestLesPresentationsSontLues:
    def test_un_pptx_est_extrait(self, tmp_path):
        """python-pptx est déjà une dépendance du projet (skills Office).

        L'accepter à l'indexation sans l'extraire n'avait donc aucune raison
        technique.
        """
        pytest = __import__("pytest")
        try:
            from pptx import Presentation
        except ImportError:  # pragma: no cover
            pytest.skip("python-pptx absent")

        presentation = Presentation()
        diapositive = presentation.slides.add_slide(presentation.slide_layouts[5])
        diapositive.shapes.title.text = "Feuille de route Synoptïa"
        chemin = tmp_path / "presentation.pptx"
        presentation.save(chemin)

        from app.services.file_parser import extract_text

        texte = extract_text(chemin)

        assert "Feuille de route Synoptïa" in texte


class TestToutesLesVuesSontAtteignablesDepuisLeChat:
    """La vue Fichiers manquait à la table de navigation du backend.

    Conséquence relevée par l'inventaire : `{action: ouvrir fichiers}` ne
    fonctionnait pas, et la réponse de `/aide` — la seule aide dérivée du code —
    ne mentionnait pas la vue qui porte l'indexation et la recherche
    documentaire. Une fonction centrale, invisible des deux seuls endroits qui
    listent ce que THÉRÈSE sait ouvrir.
    """

    def test_la_vue_fichiers_est_une_cible_de_navigation(self):
        from app.services.chat_actions import NAVIGATION_TARGETS

        assert "fichiers" in NAVIGATION_TARGETS, (
            "la vue Fichiers reste inatteignable par {action: ouvrir …} et "
            "absente de /aide, alors qu'elle porte l'indexation et le RAG"
        )

    def test_chaque_vue_de_l_application_a_une_cible(self):
        """Gate : une vue ajoutée sans cible redeviendrait invisible en silence."""
        from app.services.chat_actions import NAVIGATION_TARGETS

        # Les onze vues déclarées côté frontend (`AppView`), hors `chat` qui
        # n'est pas une destination de navigation.
        vues_attendues = {
            "accueil", "memoire", "crm", "email", "calendrier",
            "taches", "facturation", "fichiers", "projets", "documents",
        }
        manquantes = vues_attendues - set(NAVIGATION_TARGETS)

        assert not manquantes, (
            f"vues sans cible de navigation : {sorted(manquantes)} — elles "
            "seront absentes de /aide et de {action: ouvrir …}"
        )


class TestLesTroncaturesSontAnnoncees:
    """Un tableur volumineux était lu en entier, sans borne ni signalement.

    Le PDF est plafonné à 100 pages et le CSV à 500 lignes, chacun avec une
    mention de troncature dans le texte transmis. Le tableur parcourait toutes
    les feuilles et toutes les lignes : un classeur de plusieurs dizaines de
    milliers de lignes produisait un texte énorme, lent à découper et coûteux à
    envoyer au modèle, sans que rien ne l'annonce.
    """

    def test_un_gros_tableur_est_borne_et_le_dit(self, tmp_path):
        pytest = __import__("pytest")
        try:
            from openpyxl import Workbook
        except ImportError:  # pragma: no cover
            pytest.skip("openpyxl absent")

        from app.services.file_parser import MAX_XLSX_LIGNES, extract_text

        classeur = Workbook()
        feuille = classeur.active
        for numero in range(MAX_XLSX_LIGNES + 250):
            feuille.append([f"ligne-{numero}", numero])
        chemin = tmp_path / "gros-tableur.xlsx"
        classeur.save(chemin)

        texte = extract_text(chemin)

        assert "tronqué" in texte.lower(), (
            "le tableur est coupé sans que rien ne le dise : le modèle "
            "présentera une lecture partielle comme complète"
        )
        assert "ligne-0" in texte, "le début du tableur doit être transmis"
        assert f"ligne-{MAX_XLSX_LIGNES + 200}" not in texte

    def test_un_petit_tableur_passe_entier_et_sans_mention(self, tmp_path):
        """Verrou : ne pas alarmer sur un fichier lu intégralement."""
        pytest = __import__("pytest")
        try:
            from openpyxl import Workbook
        except ImportError:  # pragma: no cover
            pytest.skip("openpyxl absent")

        from app.services.file_parser import extract_text

        classeur = Workbook()
        feuille = classeur.active
        feuille.append(["Client", "Montant"])
        feuille.append(["Dupont", 1200])
        chemin = tmp_path / "petit-tableur.xlsx"
        classeur.save(chemin)

        texte = extract_text(chemin)

        assert "Dupont" in texte
        assert "tronqué" not in texte.lower()
