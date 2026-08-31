"""Replis Office : un échec franc vaut mieux qu'un livrable faux.

Revue du 30/08/2026 (soso-chemins-echec findings 1-3, grok-chemin-nominal
findings 1-2-4). Quand le code du modèle plante, THÉRÈSE livrait quand même
un fichier téléchargeable : classeur A/B/C vide, diapositives collées avec
le script Python, Word rempli de `from docx import`, page HTML d'un refus.
L'utilisateur croit que c'est fait.

Ces tests verrouillent l'inverse : sans contenu exploitable, on lève, on
ne livre pas.
"""
from __future__ import annotations

from pathlib import Path

import pytest
from app.services.skills.base import SkillParams
from app.services.skills.code_executor import (
    LivrableInexploitable,
    _validate_document_content,
)
from openpyxl import Workbook, load_workbook


def _params(title: str, content: str) -> SkillParams:
    return SkillParams(title=title, content=content)


def _xlsx_valeurs(path: Path) -> list[object]:
    wb = load_workbook(path, read_only=True)
    try:
        ws = wb.active
        assert ws is not None
        return [c.value for row in ws.iter_rows(max_row=20) for c in row if c.value]
    finally:
        wb.close()


# --- 1. XLSX : plus de classeur « Tableau / A / B / C » --------------------


PYTHON_XLSX_MERGECELL = """```python
from openpyxl import Workbook
wb = Workbook()
ws = wb.active
ws.merge_cells('A1:C1')
ws['B1'] = 'valeur'
wb.save(output_path)
```"""


@pytest.mark.asyncio
async def test_xlsx_echec_code_sans_tableau_ne_livre_pas_de_coquille(tmp_path: Path):
    """Incident 30/08 : MergedCell en écriture, puis classeur A/B/C livré."""
    from app.services.skills.xlsx_generator import XlsxSkill

    skill = XlsxSkill(output_dir=tmp_path)
    with pytest.raises(LivrableInexploitable) as caught:
        await skill.execute(_params("Suivi de chantier", PYTHON_XLSX_MERGECELL))

    assert "exploitable" in str(caught.value).lower()

    for fichier in tmp_path.glob("*.xlsx"):
        valeurs = {str(v) for v in _xlsx_valeurs(fichier)}
        assert not {"A", "B", "C"} <= valeurs, (
            f"{fichier.name} est la coquille Tableau/A/B/C livrée après l'échec"
        )


@pytest.mark.asyncio
async def test_xlsx_tableau_markdown_reel_est_accepte(tmp_path: Path):
    """Le repli reste légitime quand le modèle a vraiment écrit un tableau."""
    from app.services.skills.xlsx_generator import XlsxSkill

    skill = XlsxSkill(output_dir=tmp_path)
    markdown = (
        "## Planning\n"
        "| Jour | Heure | Evenement |\n"
        "|------|-------|----------|\n"
        "| Lundi | 10h | Visite Martin |\n"
        "| Mercredi | 14h | Fournisseur |\n"
    )
    result = await skill.execute(_params("Planning chantier", markdown))
    assert result.file_path.exists()
    valeurs = {str(v) for v in _xlsx_valeurs(result.file_path)}
    assert "Lundi" in valeurs
    assert "Visite Martin" in valeurs
    assert "A" not in valeurs


def test_xlsx_parse_sans_tableau_reste_vide():
    """Verrou : le parseur ne fabrique plus les colonnes A/B/C."""
    from app.services.skills.xlsx_generator import XlsxSkill

    skill = XlsxSkill(output_dir=Path("/tmp"))
    assert skill._parse_markdown_tables("Aucun tableau, juste une phrase.") == []


def test_xlsx_source_sans_entetes_abc_inventes():
    src = Path("src/backend/app/services/skills/xlsx_generator.py").read_text(
        encoding="utf-8"
    )
    assert '["A", "B", "C"]' not in src


def test_xlsx_garde_refuse_titre_et_entetes_abc(tmp_path: Path):
    """La garde « assez riche » comptait le titre et A/B/C comme des données."""
    path = tmp_path / "coquille.xlsx"
    wb = Workbook()
    ws = wb.active
    ws["A1"] = "Tableau"
    ws["A2"] = "A"
    ws["B2"] = "B"
    ws["C2"] = "C"
    ws["A5"] = "Généré par THERESE - Synoptia"
    wb.save(path)

    assert _validate_document_content(str(path), "xlsx") is False


# --- 2. PPTX : plus de code Python collé sur une diapositive ---------------


PYTHON_PPTX_TRIANGLE = """```python
from pptx import Presentation
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.util import Inches
prs = Presentation()
blank = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank)
slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.TRIANGLE, Inches(1), Inches(1), Inches(2), Inches(2))
prs.save(output_path)
```"""


def _pptx_textes(path: Path) -> str:
    from pptx import Presentation

    prs = Presentation(str(path))
    morceaux: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                morceaux.append(shape.text_frame.text)
    return "\n".join(morceaux)


@pytest.mark.asyncio
async def test_pptx_echec_code_sans_slides_markdown_ne_colle_pas_le_python(
    tmp_path: Path,
):
    """Incident 30/08 : MSO_AUTO_SHAPE_TYPE.TRIANGLE, le repli collait le script."""
    from app.services.skills.pptx_generator import PptxSkill

    skill = PptxSkill(output_dir=tmp_path)
    with pytest.raises(LivrableInexploitable):
        await skill.execute(_params("Kickoff client", PYTHON_PPTX_TRIANGLE))

    for fichier in tmp_path.glob("*.pptx"):
        texte = _pptx_textes(fichier)
        assert "MSO_AUTO_SHAPE_TYPE" not in texte
        assert "add_shape" not in texte
        assert "from pptx" not in texte


@pytest.mark.asyncio
async def test_pptx_markdown_reel_est_accepte(tmp_path: Path):
    from app.services.skills.pptx_generator import PptxSkill

    skill = PptxSkill(output_dir=tmp_path)
    markdown = (
        "## Contexte\n"
        "- Le chantier Martin demarre lundi\n"
        "- Budget tenu\n"
        "\n---\n\n"
        "## Planning\n"
        "- Semaine 1 : terrassement\n"
        "- Semaine 2 : fondations\n"
    )
    result = await skill.execute(_params("Kickoff Martin", markdown))
    assert result.file_path.exists()
    texte = _pptx_textes(result.file_path)
    assert "terrassement" in texte
    assert "from pptx" not in texte


# --- 3. PPTX : les diapositives ne se répètent plus ------------------------


PYTHON_PPTX_CYCLE = """```python
from pptx import Presentation
from pptx.util import Inches
prs = Presentation()
blank = prs.slide_layouts[6]
titres = ["Titre"] + ["Partie 1", "Partie 2", "Partie 3"] * 3
for titre in titres:
    slide = prs.slides.add_slide(blank)
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(10), Inches(1))
    box.text_frame.paragraphs[0].text = titre
prs.save(output_path)
```"""


@pytest.mark.asyncio
async def test_pptx_cycle_de_parties_n_est_pas_livre(tmp_path: Path):
    """Constat 30/08 : 3 parties demandées, rendues trois fois sur 10 diapos."""
    from app.services.skills.pptx_generator import PptxSkill

    skill = PptxSkill(output_dir=tmp_path)
    with pytest.raises(LivrableInexploitable):
        await skill.execute(_params("Conduite du changement", PYTHON_PPTX_CYCLE))


def test_pptx_parse_ecrase_le_cycle_repete():
    """Même cycle en Markdown de repli : on ne garde qu'un tour."""
    from app.services.skills.pptx_generator import PptxSkill

    skill = PptxSkill(output_dir=Path("/tmp"))
    bloc = (
        "## Partie 1\n- a\n\n---\n\n"
        "## Partie 2\n- b\n\n---\n\n"
        "## Partie 3\n- c\n\n---\n\n"
    )
    slides = skill._parse_content(bloc * 3)
    titres = [s["title"] for s in slides]
    assert titres == ["Partie 1", "Partie 2", "Partie 3"], titres


# --- 4. DOCX : plus de Python tronqué dans le document ---------------------


PYTHON_DOCX_TRONQUE = """Voici le document :

```python
from docx import Document
from docx.shared import Pt
doc = Document()
doc.add_heading(title, level=0)
para = doc.add_paragraph("Introduction du rapport")
run = para.add_run(
"""


@pytest.mark.asyncio
async def test_docx_python_tronque_sans_prose_ne_livre_pas(tmp_path: Path):
    """Finding Soso 2 : fence orpheline, le repli gardait les lignes Python."""
    from app.services.skills.docx_generator import DocxSkill

    skill = DocxSkill(output_dir=tmp_path)
    with pytest.raises(LivrableInexploitable):
        await skill.execute(_params("Compte-rendu client", PYTHON_DOCX_TRONQUE))

    for fichier in tmp_path.glob("*.docx"):
        from docx import Document

        texte = "\n".join(p.text for p in Document(str(fichier)).paragraphs)
        assert "from docx import" not in texte
        assert "doc.add_heading" not in texte


@pytest.mark.asyncio
async def test_docx_prose_markdown_reelle_est_acceptee(tmp_path: Path):
    from app.services.skills.docx_generator import DocxSkill

    skill = DocxSkill(output_dir=tmp_path)
    markdown = (
        "## Contexte\n\n"
        "Le client Martin valide le planning.\n\n"
        "## Suite\n\n"
        "Livraison vendredi.\n"
    )
    result = await skill.execute(_params("CR Martin", markdown))
    from docx import Document

    texte = "\n".join(p.text for p in Document(str(result.file_path)).paragraphs)
    assert "planning" in texte
    assert "from docx" not in texte


# --- 5. HTML : un refus n'est pas une page réussie -------------------------


@pytest.mark.asyncio
async def test_html_refus_du_modele_n_est_pas_une_page(tmp_path: Path):
    from app.services.skills.html_generator import HtmlSkill

    skill = HtmlSkill(output_dir=tmp_path)
    with pytest.raises(LivrableInexploitable):
        await skill.execute(
            _params(
                "Landing coach",
                "Je ne peux pas générer cette page. Voici plutôt un résumé.",
            )
        )
    assert list(tmp_path.glob("*.html")) == []


@pytest.mark.asyncio
async def test_html_partiel_sans_fermeture_n_est_pas_une_page(tmp_path: Path):
    from app.services.skills.html_generator import HtmlSkill

    skill = HtmlSkill(output_dir=tmp_path)
    with pytest.raises(LivrableInexploitable):
        await skill.execute(
            _params("Landing", "<html><body><h1>Bonjour")
        )
    assert list(tmp_path.glob("*.html")) == []


@pytest.mark.asyncio
async def test_html_complet_est_accepte(tmp_path: Path):
    from app.services.skills.html_generator import HtmlSkill

    skill = HtmlSkill(output_dir=tmp_path)
    html = (
        "<!DOCTYPE html><html lang='fr'><head><title>Coach</title></head>"
        "<body><h1>Bienvenue</h1><p>Séance d'essai.</p></body></html>"
    )
    result = await skill.execute(_params("Coach", html))
    assert result.file_path.exists()
    contenu = result.file_path.read_text(encoding="utf-8")
    assert "Bienvenue" in contenu
    assert "<pre>" not in contenu


# --- Garde « document assez riche » ----------------------------------------


def test_pptx_titre_et_merci_insuffisant(tmp_path: Path):
    """Deux diapos (titre + Merci) passaient le seuil pptx: 2."""
    from pptx import Presentation
    from pptx.util import Inches

    path = tmp_path / "coquille.pptx"
    prs = Presentation()
    blank = prs.slide_layouts[6]
    for texte in ("Suivi de chantier", "Merci"):
        slide = prs.slides.add_slide(blank)
        box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
        box.text_frame.paragraphs[0].text = texte
    prs.save(str(path))

    assert _validate_document_content(str(path), "pptx") is False


def test_garde_refuse_un_fichier_illisible(tmp_path: Path):
    """Fail-open : une exception de lecture acceptait le fichier (revue 30/08)."""
    path = tmp_path / "brise.xlsx"
    path.write_bytes(b"ceci n'est pas un classeur")
    assert _validate_document_content(str(path), "xlsx") is False


def test_pptx_une_seule_diapositive_de_fin(tmp_path: Path):
    """31/08 : le PPTX « dauphins » sortait avec deux « Merci », la seconde vide.

    Le générateur ajoute toujours sa diapositive de clôture. Quand le modèle
    en écrit une lui aussi (`## Merci`), elle est parsée comme une diapositive
    de contenu sans points, et on en livre deux.
    """
    from app.services.skills.pptx_generator import PptxSkill

    skill = PptxSkill(tmp_path)
    # Le parseur découpe sur `---`, pas sur les titres : c'est le format
    # que le prompt impose au modèle, et celui de la sortie réelle.
    contenu = (
        "## Qui sont les dauphins ?\n\n"
        "- Des cétacés à dents\n- Présents dans tous les océans\n"
        "\n---\n"
        "## Comment les protéger ?\n\n"
        "- Réduire les filets dérivants\n- Limiter le bruit sous-marin\n"
        "\n---\n"
        "## Merci\n"
    )
    slides = skill._parse_content(contenu)
    titres = [s["title"] for s in slides]
    assert "Merci" not in titres, (
        f"la clôture du modèle doit être retirée, le générateur pose la sienne : {titres}"
    )
